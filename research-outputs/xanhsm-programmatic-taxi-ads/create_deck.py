#!/usr/bin/env python3
"""
Generate PowerPoint presentation for XanhSM Programmatic Taxi Advertising

Brand: OS Research
Colors extracted from Brand Guidelines v1.0 (05/12/2025)
"""

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation with widescreen aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============== OS RESEARCH BRAND COLORS ==============
# Primary: OS Blue - #1652F0 (vibrant blue)
# Secondary: OS Blue Light - #5B8DEF
# Tertiary: OS Blue Soft - #A8C5F7
# Background: #E8F0FF (very light blue)
# Dark: #1A1A2E

OS_BLUE = RGBColor(0x16, 0x52, 0xF0)          # Primary brand color
OS_BLUE_DARK = RGBColor(0x0D, 0x3D, 0xB8)     # Darker accent
OS_BLUE_LIGHT = RGBColor(0x5B, 0x8D, 0xEF)    # Secondary blue
OS_BLUE_SOFT = RGBColor(0xA8, 0xC5, 0xF7)     # Soft accent
OS_BACKGROUND = RGBColor(0xE8, 0xF0, 0xFF)    # Light background
OS_DARK = RGBColor(0x1A, 0x1A, 0x2E)          # Dark text
OS_GRAY = RGBColor(0x6B, 0x72, 0x80)          # Body text
OS_LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)    # Subtle backgrounds
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MED_GRAY = RGBColor(0xDD, 0xDD, 0xDD)

# Alias for backward compatibility in existing functions
DARK_BLUE = OS_BLUE
ACCENT_BLUE = OS_BLUE_LIGHT
LIGHT_GRAY = OS_LIGHT_GRAY
DARK_GRAY = OS_DARK
LIGHT_BLUE = OS_BLUE_SOFT

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight_box=None):
    """Add a content slide with bullets"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Bullets
    bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(4.5))
    tf = bullet_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)
        p.level = 0

    # Highlight box if provided
    if highlight_box:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.2), Inches(12.333), Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = ACCENT_BLUE

        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(6.35), Inches(12), Inches(0.7))
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = highlight_box
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_stats_slide(prs, title, stats, subtitle_bullets=None):
    """Add a slide with stat boxes"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Stat boxes
    num_stats = len(stats)
    box_width = Inches(2.8)
    box_height = Inches(1.8)
    total_width = num_stats * box_width + (num_stats - 1) * Inches(0.3)
    start_x = (prs.slide_width - total_width) / 2
    gap = Inches(0.3)
    y_pos = Inches(1.8)

    for i, (value, label) in enumerate(stats):
        x = start_x + i * (box_width + gap)

        # Box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, box_width, box_height)
        box.fill.solid()
        box.fill.fore_color.rgb = ACCENT_BLUE
        box.line.fill.background()

        # Value
        val_box = slide.shapes.add_textbox(x, y_pos + Inches(0.25), box_width, Inches(0.8))
        tf = val_box.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Label
        lbl_box = slide.shapes.add_textbox(x, y_pos + Inches(1), box_width, Inches(0.7))
        tf = lbl_box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Subtitle bullets if provided
    if subtitle_bullets:
        bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(2.5))
        tf = bullet_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(subtitle_bullets):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(8)

    return slide

def add_comparison_slide(prs, title, headers, rows):
    """Add a comparison table slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Table
    table_width = Inches(12.333)
    col_width = table_width / len(headers)
    row_height = Inches(0.55)
    start_y = Inches(1.6)
    start_x = Inches(0.5)

    # Headers
    for i, header in enumerate(headers):
        x = start_x + i * col_width
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, start_y, col_width - Inches(0.03), row_height)
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.fill.background()

        text_box = slide.shapes.add_textbox(x + Inches(0.05), start_y + Inches(0.12), col_width - Inches(0.1), row_height)
        tf = text_box.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # Rows
    for row_idx, row in enumerate(rows):
        y = start_y + (row_idx + 1) * row_height
        for col_idx, cell in enumerate(row):
            x = start_x + col_idx * col_width

            # Alternating row colors
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, col_width - Inches(0.03), row_height)
            box.fill.solid()
            box.fill.fore_color.rgb = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
            box.line.color.rgb = MED_GRAY

            text_box = slide.shapes.add_textbox(x + Inches(0.08), y + Inches(0.12), col_width - Inches(0.15), row_height)
            tf = text_box.text_frame
            p = tf.paragraphs[0]
            p.text = cell
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            if col_idx == 0:
                p.font.bold = True

    return slide

# ============== CREATE SLIDES ==============

# Slide 1: Title
add_title_slide(prs,
    "PROGRAMMATIC TAXI SCREENS",
    "Partnering with Vietnam's #1 Electric Taxi  •  December 2025"
)

# Slide 2: The Problem
add_content_slide(prs,
    "THE PROBLEM: Advertisers Can't Measure Outdoor Ads",
    [
        "❌  No ROI measurement — \"OOH is highly impactful but notoriously difficult to measure\"",
        "❌  No targeting capabilities — Billboards reach everyone indiscriminately",
        "❌  High costs with uncertain returns — Tens to hundreds of millions VND",
        "❌  Long-term commitments required — Clients reluctant to commit",
        "❌  Limited engagement time — Only 5-7 seconds for billboard viewing",
    ],
    highlight_box="52% of marketers concerned about measurement  •  Billboard revenue fell 40% YoY"
)

# Slide 3: The Solution
add_content_slide(prs,
    "THE SOLUTION: Programmatic Taxi Screen Advertising",
    [
        "Helping advertisers reach the right people with measurable,",
        "geo-targeted ads on XanhSM's 30,000+ electric taxi screens",
        "",
        "📍  GEO-TARGET — Target by district, route, time of day",
        "📊  MEASURE — Verified impressions, GPS-tracked",
        "⚡  REAL-TIME — Launch same-day, no long-term contracts",
    ],
    highlight_box="💡 First programmatic taxi advertising platform in Vietnam"
)

# Slide 4: Market Opportunity
add_content_slide(prs,
    "MARKET OPPORTUNITY: $282M by 2033",
    [
        "Vietnam DOOH Market: $107.5M (2024) → $282.4M (2033)",
        "",
        "📈  10.14% CAGR — Highest DOOH growth rate in Southeast Asia",
        "📈  +37% programmatic DOOH volume growth (2024 vs 2023) in SEA",
        "📈  7.52% Transit Media CAGR — Fastest growing OOH segment",
        "",
        "By 2027: 75% of OOH will be DOOH, 16% programmatic",
    ]
)

# Slide 5: The Gap
add_comparison_slide(prs,
    "THE GAP: No Programmatic Taxi Advertising in Vietnam",
    ["Company", "Focus", "Programmatic?", "Gap"],
    [
        ["JCDecaux", "Airports, roadside", "✅ Yes (VIOOH)", "❌ No taxis"],
        ["Chicilon Media", "Elevator screens", "❌ Static only", "❌ No taxis"],
        ["SHOJIKI", "Taxi advertising", "❌ Static wraps", "❌ No measurement"],
        ["Taxi Ads VN", "Taxi advertising", "❌ Static decals", "❌ No targeting"],
    ]
)

# Slide 6: XanhSM Partnership
add_stats_slide(prs,
    "STRATEGIC ASSET: XanhSM — Vietnam's #1 Ride-Hailing",
    [
        ("40%", "Market Share\n#1 in Vietnam"),
        ("30,000+", "Electric Taxis\n100K total"),
        ("56/63", "Provinces\nNational Coverage"),
        ("83%", "Satisfaction\nvs Grab 80%"),
    ],
    subtitle_bullets=[
        "✅ Overtook Grab in Q1 2025 — fastest growing",
        "✅ In-car 10\" screens already installed (VinFast VF e34)",
        "✅ Premium audience — affluent urban professionals",
        "✅ No existing ad partnerships — greenfield opportunity",
    ]
)

# Slide 7: How It Works
add_content_slide(prs,
    "HOW IT WORKS: Programmatic Taxi Advertising",
    [
        "ADVERTISER → PLATFORM → TAXI SCREEN → PASSENGER",
        "",
        "1️⃣  Advertiser sets target: \"Show ads near my stores in District 1\"",
        "2️⃣  Platform matches taxis passing through target areas (GPS)",
        "3️⃣  10\" screen displays dynamic ad to captive passenger",
        "4️⃣  Dashboard shows real-time analytics and verified impressions",
        "",
        "Average ride: 8+ minutes of engagement (vs 5-7 sec billboard)",
    ]
)

# Slide 8: Target Customers
add_content_slide(prs,
    "TARGET CUSTOMERS: Who Buys OOH in Vietnam",
    [
        "OOH Spend by Segment:",
        "",
        "████████████████  FMCG/Retail — 44.5% (Vinamilk, Masan, Unilever)",
        "████████  Real Estate — 15%",
        "██████  Automotive — 10%",
        "██████  Tech/Telecom — 10%",
        "████  F&B — 8%",
        "",
        "Primary Agency Targets: GroupM, Dentsu, Publicis, Omnicom, IPG",
    ],
    highlight_box="78% of digital ad spend will be programmatic by 2028"
)

# Slide 9: Value Proposition
add_comparison_slide(prs,
    "VALUE PROPOSITION: Why Advertisers Choose Us",
    ["Feature", "Static Billboards", "Our Platform"],
    [
        ["Measurement", "❌ Traffic estimates", "✅ GPS-verified impressions"],
        ["Targeting", "❌ Everyone sees it", "✅ Geo-fence, daypart, demographics"],
        ["Commitment", "❌ 3-6 month contracts", "✅ Launch same-day"],
        ["CPM", "$5-8 (no proof)", "$8-12 (verified)"],
        ["Creative", "❌ Static image", "✅ Dynamic (weather, time, location)"],
        ["Reach", "❌ Fixed location", "✅ 30,000+ moving taxis"],
    ]
)

# Slide 10: Business Model
add_content_slide(prs,
    "BUSINESS MODEL: Revenue Sharing",
    [
        "Advertiser pays $10 CPM",
        "",
        "                    ↓",
        "",
        "    PLATFORM (60% = $6)           XanhSM (40% = $4)",
        "    • Inventory management         • Fleet access",
        "    • DSP integration              • Driver coordination",
        "    • Real-time dashboard          • Installation support",
        "",
        "Hardware: $200/vehicle (one-time) — paid by platform",
    ]
)

# Slide 11: Financial Projections
add_stats_slide(prs,
    "FINANCIAL PROJECTIONS: Path to $9.5M Revenue",
    [
        ("$1.9M", "Year 1\n1,000 taxis"),
        ("$4.5M", "Year 2\n3,000 taxis"),
        ("$9.5M", "Year 3\n5,000 taxis"),
    ],
    subtitle_bullets=[
        "Assumptions:",
        "• $10 CPM average",
        "• 8 hours/day operation, 60 ads/hour",
        "• 60/40 revenue split (platform/XanhSM)",
    ]
)

# Slide 12: Top Hypotheses
add_content_slide(prs,
    "TOP HYPOTHESES: What We Need to Prove",
    [
        "DESIRABILITY — \"Do they want this?\"",
        "  #1  FMCG brands will commit $50K+ to 3-month pilot",
        "  #2  Agencies rank \"real-time data\" as #1 must-have feature",
        "  #4  Advertisers will pay $8-12 CPM (30-50% premium)",
        "",
        "FEASIBILITY — \"Can we do this?\"",
        "  #3  XanhSM grants 1,000 taxis, ≤$200/vehicle, 90 days",
        "  #5  DSP integration (Trade Desk, DV360) feasible in 6 months",
    ]
)

# Slide 13: Experiment Plan
add_comparison_slide(prs,
    "EXPERIMENT PLAN: 9 Weeks to Validation",
    ["#", "Experiment", "Duration", "Cost", "Evidence"],
    [
        ["1", "Media Buyer Survey", "2 weeks", "$0", "Weak"],
        ["2", "Agency Interviews", "3 weeks", "$500", "Medium"],
        ["3", "XanhSM Partnership LOI", "4 weeks", "$0", "Very Strong"],
        ["4", "CPM Pricing Pre-Sale", "3 weeks", "$1,200", "Very Strong"],
        ["5", "DSP Integration Spike", "4 weeks", "$2,000", "Strong"],
    ]
)

# Slide 14: Success Criteria
add_comparison_slide(prs,
    "SUCCESS CRITERIA: Decision Gates",
    ["Experiment", "Success = GO", "Failure = PIVOT"],
    [
        ["#1 Survey", "≥70% rank real-time #1", "Pivot feature focus"],
        ["#2 Interviews", "≥50% commit $50K+", "Smaller pilot scope"],
        ["#3 XanhSM LOI", "Signed 3/4 key terms", "Approach Grab or Be"],
        ["#4 Pre-Sales", "≥$150K at $8-12 CPM", "Lower CPM / volume"],
        ["#5 DSP Spike", "Both feasible 90 days", "Self-serve first"],
    ]
)

# Slide 15: Next Steps
add_content_slide(prs,
    "NEXT STEPS: Immediate Actions",
    [
        "THIS WEEK:",
        "  ☐  Build LinkedIn contact list: 50 agency media buyers",
        "  ☐  Create 5-question survey (Google Forms)",
        "  ☐  Research XanhSM decision-makers (Head of Partnerships)",
        "  ☐  Secure warm intro to XanhSM via investor/EV network",
        "",
        "MONTH 1 TARGETS:",
        "  ☐  30+ survey responses from media buyers",
        "  ☐  12 interviews with agency/brand decision-makers",
        "  ☐  1st meeting with XanhSM partnerships team",
    ],
    highlight_box="🎯 Ultimate Success: $150K+ pre-sold + XanhSM LOI signed"
)

# Save presentation
output_path = "/Users/nntrvy/osr-company/os-research-engine/research-outputs/xanhsm-programmatic-taxi-ads/XanhSM-Programmatic-Taxi-Ads-Deck.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
